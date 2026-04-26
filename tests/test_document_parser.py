"""Tests for backend.tools.document_parser — multi-format text extraction."""
from __future__ import annotations

import io

import pytest

from backend.tools.document_parser import (
    ALLOWED_EXTENSIONS,
    MAX_COMBINED_TEXT,
    extract_text_from_file,
    extract_text_from_files,
)


# ---------------------------------------------------------------------------
# Fixture helpers — build small in-memory documents
# ---------------------------------------------------------------------------


def _make_docx(text: str) -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(text: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx(rows: list[list[str]]) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pdf(text: str) -> bytes:
    """Create a minimal single-page PDF containing *text*."""
    import pdfplumber
    # Build a trivially simple PDF by hand.
    # pdfplumber can read this back.
    from io import BytesIO

    # Use reportlab if available, otherwise build a raw PDF.
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas

        buf = BytesIO()
        c = canvas.Canvas(buf, pagesize=letter)
        c.drawString(72, 700, text)
        c.save()
        return buf.getvalue()
    except ImportError:
        # Build a minimal raw PDF manually
        content = (
            b"%PDF-1.4\n"
            b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
            b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
            b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]"
            b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
        )
        stream = f"BT /F1 12 Tf 72 700 Td ({text}) Tj ET".encode()
        stream_obj = (
            f"4 0 obj<</Length {len(stream)}>>stream\n".encode()
            + stream
            + b"\nendstream\nendobj\n"
        )
        font_obj = b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        body = content + stream_obj + font_obj
        xref_offset = len(body)
        xref = (
            b"xref\n0 6\n"
            b"0000000000 65535 f \n"
        )
        # Simplified — pdfplumber tolerates loose xref
        trailer = (
            f"trailer<</Size 6/Root 1 0 R>>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode()
        return body + xref + trailer


# ---------------------------------------------------------------------------
# Tests — extract_text_from_file
# ---------------------------------------------------------------------------


class TestExtractTextFromFile:

    def test_txt(self):
        text = extract_text_from_file(b"Hello world", "readme.txt")
        assert "Hello world" in text

    def test_html(self):
        html = b"<html><body><p>Paragraph one</p><script>ignore</script></body></html>"
        text = extract_text_from_file(html, "page.html")
        assert "Paragraph one" in text
        assert "ignore" not in text

    def test_htm_extension(self):
        html = b"<p>Content</p>"
        text = extract_text_from_file(html, "page.htm")
        assert "Content" in text

    def test_docx(self):
        data = _make_docx("Quarterly sales results")
        text = extract_text_from_file(data, "report.docx")
        assert "Quarterly sales results" in text

    def test_pptx(self):
        data = _make_pptx("Pitch deck title")
        text = extract_text_from_file(data, "deck.pptx")
        assert "Pitch deck title" in text

    def test_xlsx(self):
        data = _make_xlsx([["Company", "Revenue"], ["Acme", "10M"]])
        text = extract_text_from_file(data, "data.xlsx")
        assert "Company" in text
        assert "Acme" in text
        assert "10M" in text

    def test_pdf(self):
        data = _make_pdf("Enterprise cloud platform")
        text = extract_text_from_file(data, "brochure.pdf")
        assert "Enterprise cloud platform" in text

    def test_empty_file_returns_empty(self):
        assert extract_text_from_file(b"", "empty.txt") == ""
        assert extract_text_from_file(b"", "empty.docx") == ""

    def test_unsupported_extension_raises(self):
        with pytest.raises(ValueError, match="Unsupported file type"):
            extract_text_from_file(b"data", "malware.exe")

    def test_unsupported_zip_raises(self):
        with pytest.raises(ValueError, match="Unsupported file type"):
            extract_text_from_file(b"data", "archive.zip")

    def test_case_insensitive_extension(self):
        text = extract_text_from_file(b"hello", "FILE.TXT")
        assert "hello" in text

    def test_utf8_with_bom(self):
        content = b"\xef\xbb\xbfHello UTF-8"
        text = extract_text_from_file(content, "file.txt")
        assert "Hello UTF-8" in text


# ---------------------------------------------------------------------------
# Tests — extract_text_from_files
# ---------------------------------------------------------------------------


class TestExtractTextFromFiles:

    def test_combines_multiple(self):
        files = [
            (b"File one content", "a.txt"),
            (b"File two content", "b.txt"),
        ]
        text = extract_text_from_files(files)
        assert "File one content" in text
        assert "File two content" in text
        assert "---" in text  # separator

    def test_skips_empty(self):
        files = [
            (b"", "empty.txt"),
            (b"Real content", "real.txt"),
        ]
        text = extract_text_from_files(files)
        assert "Real content" in text
        # Only one section, no separator
        assert text.count("---") == 0

    def test_truncates_to_max(self):
        # Create content that exceeds MAX_COMBINED_TEXT
        big = b"x" * (MAX_COMBINED_TEXT + 1000)
        text = extract_text_from_files([(big, "big.txt")])
        assert len(text) <= MAX_COMBINED_TEXT

    def test_empty_list(self):
        assert extract_text_from_files([]) == ""

    def test_mixed_formats(self):
        files = [
            (b"Plain text", "notes.txt"),
            (_make_docx("Word document"), "report.docx"),
        ]
        text = extract_text_from_files(files)
        assert "Plain text" in text
        assert "Word document" in text


# ---------------------------------------------------------------------------
# Tests — constants
# ---------------------------------------------------------------------------


class TestConstants:

    def test_allowed_extensions_complete(self):
        expected = {".pdf", ".docx", ".pptx", ".xlsx", ".html", ".htm", ".txt"}
        assert ALLOWED_EXTENSIONS == expected

    def test_max_combined_text(self):
        assert MAX_COMBINED_TEXT == 120_000
