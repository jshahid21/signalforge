"""Multi-format document text extraction for seller intelligence.

Extracts plain text from PDF, DOCX, PPTX, XLSX, HTML, and TXT files.
All extracted text feeds into the same LLM extraction prompt used by
website scraping, producing identical SellerIntelligence output.
"""
from __future__ import annotations

import io
from pathlib import PurePath

from backend.tools.web_crawler import strip_html_tags

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".html", ".htm", ".txt"}
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
MAX_FILES = 5
MAX_COMBINED_TEXT = 30_000  # chars — matches website scraping truncation


def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """Extract plain text from a single file, dispatching by extension.

    Raises ``ValueError`` for unsupported extensions.
    """
    ext = PurePath(filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Accepted: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
        )

    if not file_bytes:
        return ""

    dispatch = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".xlsx": _extract_xlsx,
        ".html": _extract_html,
        ".htm": _extract_html,
        ".txt": _extract_txt,
    }
    return dispatch[ext](file_bytes)


def extract_text_from_files(files: list[tuple[bytes, str]]) -> str:
    """Extract and combine text from multiple files.

    Each element is ``(file_bytes, filename)``.  Texts are joined with
    ``\\n\\n---\\n\\n`` and the combined result is truncated to
    ``MAX_COMBINED_TEXT`` characters.
    """
    parts: list[str] = []
    for file_bytes, filename in files:
        text = extract_text_from_file(file_bytes, filename)
        if text:
            parts.append(text)
    combined = "\n\n---\n\n".join(parts)
    return combined[:MAX_COMBINED_TEXT]


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------


def _extract_pdf(data: bytes) -> str:
    import pdfplumber

    text_parts: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n\n".join(text_parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    text_parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)
    return "\n\n".join(text_parts)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    rows: list[str] = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append("\t".join(cells))
    wb.close()
    return "\n".join(rows)


def _extract_html(data: bytes) -> str:
    html = data.decode("utf-8", errors="replace")
    return strip_html_tags(html)


def _extract_txt(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")
